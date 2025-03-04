import io
import os

from fastapi import FastAPI, Response
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Импорт для генерации текста с помощью GPT-2
from transformers import GPT2LMHeadModel, GPT2Tokenizer

# Импорт для генерации изображений с помощью Stable Diffusion
from diffusers import StableDiffusionPipeline
import torch
from fastapi.middleware.cors import CORSMiddleware

app = FastAPI()


app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


@app.get("/")
def root():
    return {"message": "Hello from FastAPI!"}


class PresentationRequest(BaseModel):
    topic: str
    slides_count: int
    theme: str  # "light" или "dark"


# Инициализация модели GPT-2 для генерации текста
tokenizer = GPT2Tokenizer.from_pretrained("gpt2")
gpt2_model = GPT2LMHeadModel.from_pretrained("gpt2")

# Инициализация Stable Diffusion Pipeline для генерации изображений
device = "cuda" if torch.cuda.is_available() else "cpu"
sd_pipeline = StableDiffusionPipeline.from_pretrained(
    "CompVis/stable-diffusion-v1-4",
    torch_dtype=torch.float16 if device == "cuda" else torch.float32
)
sd_pipeline = sd_pipeline.to(device)


def generate_slide_text(topic: str, slide_number: int) -> str:
    """
    Генерирует текст для слайда, используя GPT-2.
    """
    prompt = f"Presentation topic: {topic}. Slide {slide_number}: "
    input_ids = tokenizer.encode(prompt, return_tensors="pt")
    # Генерация текста с выбором случайных вариантов (do_sample=True)
    sample_outputs = gpt2_model.generate(
        input_ids,
        do_sample=True,
        max_length=100,
        top_p=0.95,
        top_k=60,
        num_return_sequences=1
    )
    text = tokenizer.decode(sample_outputs[0], skip_special_tokens=True)
    # Извлекаем сгенерированный текст (после начального промпта)
    generated_text = text[len(prompt):].strip()
    return generated_text


def generate_slide_image(prompt: str) -> str:
    """
    Генерирует изображение для слайда, используя Stable Diffusion.
    Сохраняет изображение во временный файл и возвращает путь к файлу.
    """
    image = sd_pipeline(prompt).images[0]
    # Генерируем имя временного файла на основе хеша промпта
    image_path = f"slide_{hash(prompt)}.png"
    image.save(image_path)
    return image_path


def create_presentation(topic: str, slides_count: int, theme: str) -> io.BytesIO:
    """
    Создаёт презентацию с заданным количеством слайдов.
    На каждом слайде автоматически добавляется текст (сгенерированный GPT-2)
    и изображение (сгенерированное Stable Diffusion). Фон слайда выбирается в зависимости от темы.
    """
    prs = Presentation()

    # Определяем цвета в зависимости от выбранной темы
    if theme.lower() == "dark":
        bg_color = RGBColor(0, 0, 0)  # тёмный фон
        text_color = RGBColor(255, 255, 255)  # белый текст
    else:
        bg_color = RGBColor(255, 255, 255)  # светлый фон
        text_color = RGBColor(0, 0, 0)  # чёрный текст

    for i in range(1, slides_count + 1):
        # Добавляем пустой слайд (layout index 5 – пустой макет, можно изменить по необходимости)
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        # Задаём цвет фона слайда
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        # Генерация текста слайда
        slide_text = generate_slide_text(topic, i)
        # Промпт для генерации изображения – можно расширять описание
        image_prompt = f"{topic} slide {i} illustration"
        image_path = generate_slide_image(image_prompt)


        # Добавляем текстовое поле
        left = Inches(1)
        top = Inches(0.5)
        width = Inches(8)
        height = Inches(1.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        p = tf.add_paragraph()
        p.text = slide_text
        p.font.size = Pt(24)
        p.font.color.rgb = text_color

        # Добавляем изображение на слайд
        pic_left = Inches(1)
        pic_top = Inches(2.5)
        pic_width = Inches(6)
        slide.shapes.add_picture(image_path, pic_left, pic_top, width=pic_width)

        # Удаляем временный файл с изображением
        if os.path.exists(image_path):
            os.remove(image_path)

    # Сохраняем презентацию в BytesIO
    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io


@app.options("/generate")
async def options_generate():
    response = Response()
    response.headers["Allow"] = "POST, OPTIONS"
    response.status_code = 200
    return response


@app.post("/generate")
async def generate_presentation(pres_req: PresentationRequest):
    """
    Endpoint, который принимает данные презентации, генерирует pptx файл
    и возвращает его в качестве ответа.
    """
    ppt_io = create_presentation(pres_req.topic, pres_req.slides_count, pres_req.theme)
    return StreamingResponse(
        ppt_io,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=presentation.pptx"}
    )


if __name__ == "__main__":
    import uvicorn

    uvicorn.run(app, host="0.0.0.0", port=8000)