import io
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pydantic import BaseModel
from text_generator import generate_slide_text
from image_generaor import generate_slide_image



class PresentationRequest(BaseModel):
    topic: str
    slides_count: int
    theme: str


def create_presentation(topic: str, slides_count: int, theme: str) -> io.BytesIO:
    """
    Создаёт презентацию с заданным количеством слайдов, используя шаблон из папки themes.
    В зависимости от темы выбирается dark.pptx или light.pptx.
    На каждом слайде добавляется текст (сгенерированный GPT-2) и изображение (сгенерированное Stable Diffusion)
    по следующему алгоритму: текст добавляется в левую часть слайда, а изображение – в правую.
    """
    # Определяем путь к шаблону
    theme = theme.lower()
    if theme == "dark":
        template_path = "themes/dark.pptx"
    else:
        template_path = "themes/light.pptx"

    # Загружаем презентацию из шаблона
    prs = Presentation(template_path)

    # Создание титульного листа
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    textbox = slide.shapes.add_textbox(
        left=Inches(1.5),
        top=Inches(2.2),
        width=Inches(10.32),
        height=Inches(2)
    )
    tf = textbox.text_frame
    tf.auto_size = False
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)
    p = tf.add_paragraph()
    p.text = topic
    p.font.size = Pt(40)
    p.alignment = PP_ALIGN.CENTER


    for i in range(1, slides_count + 1):
        # Создаем слайд на основе макета шаблона
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # Генерация текста для слайда
        slide_text = generate_slide_text(topic)
        # Промпт для генерации изображения
        image_prompt = f"{topic} slide {i} illustration"
        image_path = generate_slide_image(image_prompt)

        # Добавляем текстовое поле
        textbox = slide.shapes.add_textbox(
            left=Inches(0.5),
            top=Inches(0.5),
            width=Inches(5.5),
            height=Inches(6)
        )

        # Получаем текстовый фрейм и настраиваем его свойства
        tf = textbox.text_frame
        tf.auto_size = False
        tf.word_wrap = True

        # Устанавливаем фиксированные отступы
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0.2)

        # Добавляем абзац с текстом
        p = tf.add_paragraph()
        p.text = slide_text
        p.font.size = Pt(24)


        # Добавляем изображение
        pic_left = Inches(6.16)
        pic_top = Inches(0.6)
        pic_width = Inches(6.29)
        slide.shapes.add_picture(image_path, pic_left, pic_top, width=pic_width)

        # Удаляем временный файл с изображением, если он существует
        if os.path.exists(image_path):
            os.remove(image_path)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    textbox = slide.shapes.add_textbox(
        left=Inches(1.5),
        top=Inches(2.2),
        width=Inches(10.32),
        height=Inches(2)
    )
    tf = textbox.text_frame
    tf.auto_size = False
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)
    p = tf.add_paragraph()
    p.text = "Спасибо за внимание!"
    p.font.size = Pt(40)
    p.alignment = PP_ALIGN.CENTER

    # Сохраняем презентацию
    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io
